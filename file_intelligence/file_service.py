from __future__ import annotations

import base64
import hashlib
import io
import json
import logging
import os
import shutil
import subprocess
import tarfile
import tempfile
import time
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

import requests
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel, Field

HOST = os.getenv("AURORAFOX_FILES_HOST", "127.0.0.1")
PORT = int(os.getenv("AURORAFOX_FILES_PORT", "8767"))
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://127.0.0.1:11434").rstrip("/")
VISION_MODEL = os.getenv("AURORAFOX_VISION_MODEL", "qwen3-vl:8b")
VOICE_URL = os.getenv("AURORAFOX_VOICE_URL", "http://127.0.0.1:8765").rstrip("/")
USER_ROOT = Path(os.getenv("AURORAFOX_USER_DIR", str(Path.home() / ".aurorafox"))).resolve()
CACHE_DIR = USER_ROOT / "file_cache"
LOG_DIR = USER_ROOT / "logs"
CACHE_DIR.mkdir(parents=True, exist_ok=True)
LOG_DIR.mkdir(parents=True, exist_ok=True)

MAX_FILE_BYTES = int(os.getenv("AURORAFOX_FILE_MAX_BYTES", str(1024 * 1024 * 1024)))
MAX_TEXT_CHARS = int(os.getenv("AURORAFOX_FILE_MAX_TEXT", "160000"))
MAX_ARCHIVE_ENTRIES = int(os.getenv("AURORAFOX_ARCHIVE_MAX_ENTRIES", "5000"))
MAX_ARCHIVE_EXPANDED = int(os.getenv("AURORAFOX_ARCHIVE_MAX_EXPANDED", str(512 * 1024 * 1024)))
MAX_TREE_ITEMS = 5000

logging.basicConfig(
    filename=LOG_DIR / "aurora_files.log",
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(message)s",
    encoding="utf-8",
)
log = logging.getLogger("aurora_files")

app = FastAPI(title="AuroraFox File Intelligence", version="1.0.0")


class AnalyzeRequest(BaseModel):
    path: str = Field(min_length=1, max_length=8192)
    question: str = Field(default="", max_length=12000)
    visual: bool = True
    max_chars: int = Field(default=MAX_TEXT_CHARS, ge=2000, le=500000)


class TreeRequest(BaseModel):
    path: str = Field(min_length=1, max_length=8192)
    max_items: int = Field(default=2000, ge=1, le=MAX_TREE_ITEMS)


class CacheSearchRequest(BaseModel):
    query: str = Field(min_length=1, max_length=1000)
    limit: int = Field(default=20, ge=1, le=100)


TEXT_EXT = {
    ".txt", ".md", ".json", ".csv", ".tsv", ".gd", ".py", ".js", ".ts", ".tsx", ".jsx",
    ".html", ".css", ".scss", ".xml", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".log",
    ".shader", ".glsl", ".cpp", ".c", ".h", ".hpp", ".cs", ".java", ".kt", ".rs", ".go",
    ".php", ".rb", ".lua", ".swift", ".dart", ".sql", ".sh", ".ps1", ".r", ".jl", ".ex", ".exs",
}
IMAGE_EXT = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff"}
AUDIO_EXT = {".wav", ".mp3", ".ogg", ".flac", ".m4a", ".aac", ".opus"}
VIDEO_EXT = {".mp4", ".mkv", ".webm", ".mov", ".avi", ".m4v"}
ARCHIVE_EXT = {".zip", ".7z", ".tar", ".gz", ".tgz", ".bz2", ".tbz2", ".xz", ".txz"}


def _safe_file(path: str) -> Path:
    try:
        p = Path(path).expanduser().resolve(strict=True)
    except Exception as exc:
        raise HTTPException(404, f"File not found: {path}") from exc
    if not p.is_file():
        raise HTTPException(400, "Path is not a file")
    size = p.stat().st_size
    if size > MAX_FILE_BYTES:
        raise HTTPException(413, f"File is too large: {size} bytes")
    return p


def _safe_dir(path: str) -> Path:
    try:
        p = Path(path).expanduser().resolve(strict=True)
    except Exception as exc:
        raise HTTPException(404, f"Directory not found: {path}") from exc
    if not p.is_dir():
        raise HTTPException(400, "Path is not a directory")
    return p


def _cache_key(path: Path, question: str, visual: bool) -> str:
    st = path.stat()
    raw = f"{path}|{st.st_size}|{st.st_mtime_ns}|{question}|{visual}|v1"
    return hashlib.sha256(raw.encode("utf-8", errors="replace")).hexdigest()


def _cache_get(key: str) -> dict[str, Any] | None:
    p = CACHE_DIR / f"{key}.json"
    if not p.is_file():
        return None
    try:
        data = json.loads(p.read_text(encoding="utf-8"))
        p.touch()
        return data if isinstance(data, dict) else None
    except Exception:
        return None


def _cache_put(key: str, payload: dict[str, Any]) -> None:
    target = CACHE_DIR / f"{key}.json"
    target.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    _trim_cache(512 * 1024 * 1024)


def _trim_cache(limit: int) -> None:
    files = sorted(CACHE_DIR.glob("*.json"), key=lambda p: p.stat().st_mtime)
    total = sum(p.stat().st_size for p in files)
    while files and total > limit:
        p = files.pop(0)
        total -= p.stat().st_size
        p.unlink(missing_ok=True)


def _truncate(text: str, limit: int) -> tuple[str, bool]:
    if len(text) <= limit:
        return text, False
    return text[:limit] + "\n\n[Обрезано AuroraFox: достигнут лимит контекста]", True


def _read_text(path: Path) -> tuple[str, str]:
    raw = path.read_bytes()
    for enc in ("utf-8-sig", "utf-8", "cp1251", "utf-16", "latin-1"):
        try:
            return raw.decode(enc), enc
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace"), "utf-8-replace"


def _text_from_docx(path: Path) -> tuple[str, dict[str, Any]]:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    table_count = 0
    for table in doc.tables:
        table_count += 1
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts), {"paragraphs": len(doc.paragraphs), "tables": table_count}


def _text_from_xlsx(path: Path) -> tuple[str, dict[str, Any]]:
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    parts: list[str] = []
    sheets: list[dict[str, Any]] = []
    cells_seen = 0
    try:
        for ws in wb.worksheets:
            parts.append(f"\n### Лист: {ws.title}")
            rows_seen = 0
            for row in ws.iter_rows(values_only=True):
                values = ["" if v is None else str(v) for v in row]
                if any(v for v in values):
                    parts.append("\t".join(values))
                    rows_seen += 1
                    cells_seen += len(values)
                if cells_seen >= 50000:
                    parts.append("[Таблица обрезана: лимит 50000 ячеек]")
                    break
            sheets.append({"name": ws.title, "rows_read": rows_seen})
            if cells_seen >= 50000:
                break
    finally:
        wb.close()
    return "\n".join(parts), {"sheets": sheets, "cells_read": cells_seen}


def _text_from_xls(path: Path) -> tuple[str, dict[str, Any]]:
    import xlrd

    book = xlrd.open_workbook(str(path), on_demand=True)
    parts: list[str] = []
    sheets: list[dict[str, Any]] = []
    cells = 0
    try:
        for sheet in book.sheets():
            parts.append(f"\n### Лист: {sheet.name}")
            max_rows = min(sheet.nrows, 10000)
            for r in range(max_rows):
                values = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
                if any(v for v in values):
                    parts.append("\t".join(values))
                    cells += len(values)
                if cells >= 50000:
                    parts.append("[Таблица обрезана: лимит 50000 ячеек]")
                    break
            sheets.append({"name": sheet.name, "rows": sheet.nrows, "cols": sheet.ncols})
            if cells >= 50000:
                break
    finally:
        book.release_resources()
    return "\n".join(parts), {"sheets": sheets, "cells_read": cells}


def _text_from_pptx(path: Path) -> tuple[str, dict[str, Any]]:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                slide_parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    slide_parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        if slide_parts:
            parts.append(f"\n### Слайд {idx}\n" + "\n".join(slide_parts))
    return "\n".join(parts), {"slides": len(prs.slides)}


def _text_from_open_document(path: Path) -> tuple[str, dict[str, Any]]:
    if not zipfile.is_zipfile(path):
        raise ValueError("Invalid OpenDocument container")
    with zipfile.ZipFile(path) as zf:
        raw = zf.read("content.xml")
    root = ET.fromstring(raw)
    chunks: list[str] = []
    for elem in root.iter():
        if elem.text and elem.text.strip():
            chunks.append(elem.text.strip())
    return "\n".join(chunks), {"xml_nodes": sum(1 for _ in root.iter())}


def _pdf_extract(path: Path, visual: bool, question: str) -> tuple[str, dict[str, Any], list[str]]:
    from pypdf import PdfReader

    warnings: list[str] = []
    reader = PdfReader(str(path))
    parts: list[str] = []
    page_lengths: list[int] = []
    for idx, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception as exc:
            text = ""
            warnings.append(f"Страница {idx + 1}: ошибка извлечения текста: {exc}")
        page_lengths.append(len(text.strip()))
        if text.strip():
            parts.append(f"\n### Страница {idx + 1}\n{text}")
    text = "\n".join(parts)
    low_text = len(text.strip()) < max(250, len(reader.pages) * 80)
    if visual and low_text and reader.pages:
        try:
            rendered = _render_pdf_pages(path, min(6, len(reader.pages)))
            visual_texts = []
            for page_no, png in rendered:
                prompt = question.strip() or "Распознай текст на странице и кратко опиши важное содержимое. Ответь по-русски."
                result = _vision_bytes(png, prompt)
                if result:
                    visual_texts.append(f"\n### Визуальный анализ страницы {page_no}\n{result}")
            if visual_texts:
                text = text + "\n" + "\n".join(visual_texts)
                warnings.append("PDF содержал мало извлекаемого текста; добавлен локальный визуальный анализ выбранных страниц.")
        except Exception as exc:
            warnings.append(f"Визуальный fallback PDF недоступен: {exc}")
    return text, {"pages": len(reader.pages), "page_text_lengths": page_lengths}, warnings


def _render_pdf_pages(path: Path, count: int) -> list[tuple[int, bytes]]:
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(str(path))
    out: list[tuple[int, bytes]] = []
    try:
        for idx in range(min(count, len(pdf))):
            page = pdf[idx]
            bitmap = page.render(scale=1.4)
            pil = bitmap.to_pil()
            buf = io.BytesIO()
            pil.save(buf, format="PNG")
            out.append((idx + 1, buf.getvalue()))
            bitmap.close()
            page.close()
    finally:
        pdf.close()
    return out


def _vision_bytes(data: bytes, prompt: str) -> str:
    payload = {
        "model": VISION_MODEL,
        "stream": False,
        "messages": [{"role": "user", "content": prompt, "images": [base64.b64encode(data).decode("ascii")]}],
        "options": {"temperature": 0.1},
    }
    r = requests.post(f"{OLLAMA_URL}/api/chat", json=payload, timeout=180)
    if r.status_code != 200:
        raise RuntimeError(f"Ollama HTTP {r.status_code}: {r.text[:500]}")
    return str(r.json().get("message", {}).get("content", "")).strip()


def _image_analyze(path: Path, question: str, visual: bool) -> tuple[str, dict[str, Any], list[str]]:
    from PIL import Image

    warnings: list[str] = []
    with Image.open(path) as im:
        meta = {"width": im.width, "height": im.height, "mode": im.mode, "format": im.format, "frames": getattr(im, "n_frames", 1)}
        if not visual:
            return "Изображение: {width}×{height}, формат {format}, режим {mode}.".format(**meta), meta, warnings
        frame = im.copy()
        if frame.mode not in ("RGB", "RGBA"):
            frame = frame.convert("RGB")
        frame.thumbnail((2048, 2048))
        buf = io.BytesIO()
        frame.save(buf, format="PNG")
    prompt = question.strip() or "Подробно проанализируй изображение: прочитай видимый текст, опиши важные элементы и ответь по-русски."
    try:
        text = _vision_bytes(buf.getvalue(), prompt)
    except Exception as exc:
        text = ""
        warnings.append(f"Vision-анализ недоступен: {exc}")
    return text or f"Изображение {meta['width']}×{meta['height']}", meta, warnings


def _voice_transcribe(path: Path) -> tuple[str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    try:
        r = requests.post(f"{VOICE_URL}/stt_path", json={"path": str(path)}, timeout=300)
        if r.status_code == 200:
            data = r.json()
            if data.get("ok"):
                return str(data.get("text", "")), {"engine": "AuroraVoice"}, warnings
        warnings.append(f"Voice backend STT unavailable: HTTP {r.status_code}")
    except Exception as exc:
        warnings.append(f"Voice backend STT unavailable: {exc}")
    return "", {}, warnings


def _video_analyze(path: Path, question: str, visual: bool) -> tuple[str, dict[str, Any], list[str]]:
    import imageio_ffmpeg

    warnings: list[str] = []
    ffmpeg = imageio_ffmpeg.get_ffmpeg_exe()
    parts: list[str] = []
    frame_results: list[str] = []
    with tempfile.TemporaryDirectory(prefix="aurorafox-video-") as tmp:
        root = Path(tmp)
        audio = root / "aurorafox_voice_input.wav"
        cmd = [ffmpeg, "-y", "-i", str(path), "-vn", "-ac", "1", "-ar", "16000", str(audio)]
        proc = subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.PIPE, timeout=240, shell=False)
        if proc.returncode == 0 and audio.is_file() and audio.stat().st_size > 44:
            transcript, _, w = _voice_transcribe(audio)
            warnings.extend(w)
            if transcript:
                parts.append("### Расшифровка аудио\n" + transcript)
        else:
            warnings.append("Не удалось извлечь аудиодорожку из видео.")

        if visual:
            pattern = str(root / "frame-%02d.jpg")
            frame_cmd = [
                ffmpeg, "-y", "-i", str(path), "-vf",
                "fps=1/30,scale='min(1280,iw)':-2", "-frames:v", "8", pattern,
            ]
            proc = subprocess.run(frame_cmd, stdout=subprocess.DEVNULL, stderr=subprocess.PIPE, timeout=240, shell=False)
            if proc.returncode == 0:
                prompt = question.strip() or "Опиши, что происходит на этом кадре видео, и прочитай важный видимый текст. Ответь по-русски."
                for frame in sorted(root.glob("frame-*.jpg"))[:8]:
                    try:
                        result = _vision_bytes(frame.read_bytes(), prompt)
                        if result:
                            frame_results.append(f"{frame.stem}: {result}")
                    except Exception as exc:
                        warnings.append(f"Vision-анализ кадра недоступен: {exc}")
                        break
    if frame_results:
        parts.append("### Выбранные кадры\n" + "\n".join(frame_results))
    return "\n\n".join(parts), {"frames_analyzed": len(frame_results)}, warnings


def _archive_listing(path: Path) -> tuple[str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    entries: list[dict[str, Any]] = []
    total = 0

    def add(name: str, size: int, is_dir: bool = False) -> None:
        nonlocal total
        normalized = name.replace("\\", "/")
        p = Path(normalized)
        unsafe = p.is_absolute() or ".." in p.parts
        entries.append({"path": normalized, "size": int(size), "dir": is_dir, "unsafe": unsafe})
        if not is_dir:
            total += max(0, int(size))

    suffix = path.suffix.lower()
    if suffix == ".zip":
        with zipfile.ZipFile(path) as zf:
            for info in zf.infolist()[:MAX_ARCHIVE_ENTRIES + 1]:
                add(info.filename, info.file_size, info.is_dir())
    elif suffix == ".7z":
        import py7zr
        with py7zr.SevenZipFile(path, mode="r") as zf:
            for info in zf.list()[:MAX_ARCHIVE_ENTRIES + 1]:
                add(str(info.filename), int(getattr(info, "uncompressed", 0) or 0), bool(getattr(info, "is_directory", False)))
    elif tarfile.is_tarfile(path):
        with tarfile.open(path, mode="r:*") as tf:
            for info in tf.getmembers()[:MAX_ARCHIVE_ENTRIES + 1]:
                add(info.name, info.size, info.isdir())
    else:
        raise ValueError("Формат архива не поддерживается безопасным локальным обработчиком")

    if len(entries) > MAX_ARCHIVE_ENTRIES:
        warnings.append(f"Архив содержит больше {MAX_ARCHIVE_ENTRIES} записей; список обрезан.")
        entries = entries[:MAX_ARCHIVE_ENTRIES]
    unsafe_count = sum(1 for e in entries if e["unsafe"])
    if unsafe_count:
        warnings.append(f"Обнаружено потенциально небезопасных путей: {unsafe_count}; распаковка таких путей запрещена.")
    if total > MAX_ARCHIVE_EXPANDED:
        warnings.append(f"Заявленный распакованный размер превышает лимит {MAX_ARCHIVE_EXPANDED} байт; автоматическая распаковка запрещена.")
    text_lines = [f"{('[DIR] ' if e['dir'] else '')}{e['path']} ({e['size']} B){' [UNSAFE]' if e['unsafe'] else ''}" for e in entries]
    return "\n".join(text_lines), {"entries": len(entries), "expanded_bytes": total, "unsafe_entries": unsafe_count}, warnings


def _analyze(path: Path, question: str, visual: bool) -> dict[str, Any]:
    ext = path.suffix.lower()
    warnings: list[str] = []
    metadata: dict[str, Any] = {"name": path.name, "extension": ext, "size": path.stat().st_size}
    text = ""
    kind = "binary"

    if ext in TEXT_EXT:
        kind = "text/code"
        text, encoding = _read_text(path)
        metadata["encoding"] = encoding
    elif ext == ".pdf":
        kind = "pdf"
        text, extra, warnings = _pdf_extract(path, visual, question)
        metadata.update(extra)
    elif ext == ".docx":
        kind = "document"
        text, extra = _text_from_docx(path)
        metadata.update(extra)
    elif ext == ".xlsx":
        kind = "spreadsheet"
        text, extra = _text_from_xlsx(path)
        metadata.update(extra)
    elif ext == ".xls":
        kind = "spreadsheet"
        text, extra = _text_from_xls(path)
        metadata.update(extra)
    elif ext == ".pptx":
        kind = "presentation"
        text, extra = _text_from_pptx(path)
        metadata.update(extra)
    elif ext in {".odt", ".ods"}:
        kind = "document" if ext == ".odt" else "spreadsheet"
        text, extra = _text_from_open_document(path)
        metadata.update(extra)
    elif ext in IMAGE_EXT:
        kind = "image"
        text, extra, warnings = _image_analyze(path, question, visual)
        metadata.update(extra)
    elif ext in AUDIO_EXT:
        kind = "audio"
        text, extra, warnings = _voice_transcribe(path)
        metadata.update(extra)
    elif ext in VIDEO_EXT:
        kind = "video"
        text, extra, warnings = _video_analyze(path, question, visual)
        metadata.update(extra)
    elif ext in ARCHIVE_EXT or zipfile.is_zipfile(path) or tarfile.is_tarfile(path):
        kind = "archive"
        text, extra, warnings = _archive_listing(path)
        metadata.update(extra)
    elif ext == ".rar":
        kind = "archive"
        warnings.append("RAR принят, но автоматическая распаковка отключена: в локальный runtime не добавлен отдельный RAR backend.")
        text = "RAR-архив. Можно сохранить и обработать после подключения совместимого локального распаковщика."
    else:
        try:
            text, encoding = _read_text(path)
            if "\x00" not in text[:4096]:
                kind = "text"
                metadata["encoding"] = encoding
            else:
                text = "Бинарный файл: содержимое не преобразовано в текст."
        except Exception:
            text = "Бинарный файл: содержимое не преобразовано в текст."

    return {"kind": kind, "text": text, "metadata": metadata, "warnings": warnings}


@app.get("/health")
def health() -> dict[str, Any]:
    vision = False
    voice = False
    try:
        vision = requests.get(f"{OLLAMA_URL}/api/tags", timeout=1.5).status_code == 200
    except Exception:
        pass
    try:
        voice = requests.get(f"{VOICE_URL}/health", timeout=1.5).status_code == 200
    except Exception:
        pass
    return {
        "ok": True,
        "backend": "AuroraFileIntelligence",
        "vision_online": vision,
        "voice_online": voice,
        "cache_dir": str(CACHE_DIR),
        "limits": {"max_file_bytes": MAX_FILE_BYTES, "max_text_chars": MAX_TEXT_CHARS},
    }


@app.post("/analyze")
def analyze(req: AnalyzeRequest) -> dict[str, Any]:
    path = _safe_file(req.path)
    key = _cache_key(path, req.question, req.visual)
    cached = _cache_get(key)
    if cached is not None:
        cached["cached"] = True
        return cached
    started = time.time()
    try:
        result = _analyze(path, req.question, req.visual)
        text, truncated = _truncate(str(result.get("text", "")), req.max_chars)
        payload = {
            "ok": True,
            "path": str(path),
            "name": path.name,
            "kind": result.get("kind", "unknown"),
            "content": text,
            "metadata": result.get("metadata", {}),
            "warnings": result.get("warnings", []),
            "truncated": truncated,
            "cached": False,
            "elapsed_ms": int((time.time() - started) * 1000),
        }
        _cache_put(key, payload)
        log.info("analyzed name=%s kind=%s size=%d ms=%d", path.name, payload["kind"], path.stat().st_size, payload["elapsed_ms"])
        return payload
    except HTTPException:
        raise
    except Exception as exc:
        log.exception("analysis failed path=%s", path)
        raise HTTPException(422, f"File analysis failed: {exc}") from exc


@app.post("/tree")
def tree(req: TreeRequest) -> dict[str, Any]:
    root = _safe_dir(req.path)
    items: list[dict[str, Any]] = []
    for p in root.rglob("*"):
        if len(items) >= req.max_items:
            break
        try:
            items.append({
                "path": p.relative_to(root).as_posix(),
                "dir": p.is_dir(),
                "size": p.stat().st_size if p.is_file() else 0,
            })
        except OSError:
            continue
    return {"ok": True, "root": str(root), "items": items, "truncated": len(items) >= req.max_items}


@app.post("/cache/search")
def cache_search(req: CacheSearchRequest) -> dict[str, Any]:
    q = req.query.casefold()
    results: list[dict[str, Any]] = []
    for p in sorted(CACHE_DIR.glob("*.json"), key=lambda x: x.stat().st_mtime, reverse=True):
        try:
            data = json.loads(p.read_text(encoding="utf-8"))
        except Exception:
            continue
        hay = (str(data.get("name", "")) + "\n" + str(data.get("content", ""))).casefold()
        if q in hay:
            results.append({
                "name": data.get("name", ""),
                "path": data.get("path", ""),
                "kind": data.get("kind", ""),
                "excerpt": str(data.get("content", ""))[:1200],
            })
        if len(results) >= req.limit:
            break
    return {"ok": True, "results": results}


@app.post("/cache/clear")
def clear_cache() -> dict[str, Any]:
    removed = 0
    for p in CACHE_DIR.glob("*.json"):
        p.unlink(missing_ok=True)
        removed += 1
    return {"ok": True, "removed": removed}


@app.post("/shutdown")
def shutdown() -> dict[str, Any]:
    import threading

    def die() -> None:
        time.sleep(0.15)
        os._exit(0)

    threading.Thread(target=die, daemon=True).start()
    return {"ok": True}


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host=HOST, port=PORT, log_level="warning")
