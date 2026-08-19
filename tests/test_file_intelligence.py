from __future__ import annotations

import sys
import zipfile
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "file_intelligence"))

import file_service  # noqa: E402
from file_service import _analyze, _archive_listing  # noqa: E402


def test_text_file(tmp_path: Path) -> None:
    path = tmp_path / "hello.txt"
    path.write_text("Привет, AuroraFox!", encoding="utf-8")
    result = _analyze(path, "", False)
    assert result["kind"] == "text/code"
    assert "AuroraFox" in result["text"]
    assert result["metadata"]["encoding"].startswith("utf-8")


def test_docx_xlsx_pptx(tmp_path: Path) -> None:
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation

    docx_path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_paragraph("Документ AuroraFox")
    table = doc.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "ключ"
    table.cell(0, 1).text = "значение"
    doc.save(docx_path)
    doc_result = _analyze(docx_path, "", False)
    assert "Документ AuroraFox" in doc_result["text"]
    assert "ключ | значение" in doc_result["text"]

    xlsx_path = tmp_path / "sample.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Данные"
    ws.append(["Имя", "Значение"])
    ws.append(["AuroraFox", 42])
    wb.save(xlsx_path)
    wb.close()
    xlsx_result = _analyze(xlsx_path, "", False)
    assert "AuroraFox" in xlsx_result["text"]
    assert "42" in xlsx_result["text"]

    pptx_path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AuroraFox"
    slide.placeholders[1].text = "Локальная презентация"
    prs.save(pptx_path)
    pptx_result = _analyze(pptx_path, "", False)
    assert "AuroraFox" in pptx_result["text"]
    assert "Локальная презентация" in pptx_result["text"]


def test_zip_path_traversal_is_detected(tmp_path: Path) -> None:
    archive = tmp_path / "unsafe.zip"
    with zipfile.ZipFile(archive, "w") as zf:
        zf.writestr("safe/readme.txt", "ok")
        zf.writestr("../escape.txt", "no")

    text, meta, warnings = _archive_listing(archive)
    assert meta["entries"] == 2
    assert meta["unsafe_entries"] == 1
    assert "[UNSAFE]" in text
    assert any("небезопас" in warning.lower() for warning in warnings)


def test_binary_does_not_crash(tmp_path: Path) -> None:
    path = tmp_path / "blob.dat"
    path.write_bytes(b"\x00\x01\x02\xff" * 100)
    result = _analyze(path, "", False)
    assert result["kind"] == "binary"
    assert result["text"]


class _FakeResponse:
    def __init__(self, status_code: int, payload: dict | None = None) -> None:
        self.status_code = status_code
        self._payload = payload or {}

    def json(self) -> dict:
        return self._payload


def test_health_does_not_claim_vision_without_model(monkeypatch: pytest.MonkeyPatch) -> None:
    def fake_get(url: str, timeout: float):
        if url.endswith("/api/tags"):
            return _FakeResponse(200, {"models": [{"name": "qwen3:8b"}]})
        return _FakeResponse(200, {"ok": True})

    monkeypatch.setattr(file_service.requests, "get", fake_get)
    result = file_service.health()
    assert result["ok"] is True
    assert result["ollama_online"] is True
    assert result["vision_online"] is False
    assert result["vision_model"] == file_service.VISION_MODEL


def test_health_claims_vision_only_with_selected_model(monkeypatch: pytest.MonkeyPatch) -> None:
    def fake_get(url: str, timeout: float):
        if url.endswith("/api/tags"):
            return _FakeResponse(200, {"models": [{"name": "qwen3:8b"}, {"name": file_service.VISION_MODEL}]})
        return _FakeResponse(200, {"ok": True})

    monkeypatch.setattr(file_service.requests, "get", fake_get)
    result = file_service.health()
    assert result["ollama_online"] is True
    assert result["vision_online"] is True
    assert file_service.VISION_MODEL in result["installed_models"]
